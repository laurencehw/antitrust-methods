"""
deck_builder.py — reusable PowerPoint builder for the "Antitrust Methods" class.

Produces clean, editable .pptx decks styled to match the book's palette.
No R / Quarto required; figures are described or shown as code/data so the
deck stands alone (the live R demos are run from the book repo in class).

Usage:
    from deck_builder import Deck
    d = Deck("Day 1 — Foundations", "Antitrust Methods for Practitioners")
    d.section("Module 1.1 — The competition-law landscape")
    d.bullets("Why institutions come first", ["...", "..."])
    d.save("slides/day1_foundations.pptx")

Palette (book helpers.R):
    blue #0072B2, orange #E69F00, green #009E73, red #D55E00,
    purple #CC79A7, cyan #56B4E9, grey #555555
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette -------------------------------------------------------------
BLUE = RGBColor(0x00, 0x72, 0xB2)
ORANGE = RGBColor(0xE6, 0x9F, 0x00)
GREEN = RGBColor(0x00, 0x9E, 0x73)
RED = RGBColor(0xD5, 0x5E, 0x00)
PURPLE = RGBColor(0xCC, 0x79, 0xA7)
CYAN = RGBColor(0x56, 0xB4, 0xE9)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# jurisdiction tags
JURIS = {"US": BLUE, "EU": ORANGE, "SA": GREEN, "UK": PURPLE, "OTHER": GREY}

FONT = "Calibri"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)  # 16:9


class Deck:
    def __init__(self, title, subtitle=""):
        self.prs = Presentation()
        self.prs.slide_width = W
        self.prs.slide_height = H
        self.blank = self.prs.slide_layouts[6]
        self._title_slide(title, subtitle)

    # ---- low-level helpers ----------------------------------------------
    def _slide(self):
        return self.prs.slides.add_slide(self.blank)

    def _box(self, slide, l, t, w, h, fill=None, line=None):
        shp = slide.shapes.add_shape(1, l, t, w, h)  # rectangle
        shp.shadow.inherit = False
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(1)
        return shp

    def _text(self, slide, l, t, w, h, runs, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, space_after=6, line_spacing=1.0):
        """runs: list of paragraphs; each paragraph = list of (text, opts) tuples
        or a plain string. opts: dict(size,bold,color,italic,font,bullet,level)."""
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        first = True
        for para in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = Pt(space_after)
            p.line_spacing = line_spacing
            if isinstance(para, str):
                para = [(para, {})]
            level = 0
            for seg, opts in para:
                level = opts.get("level", level)
                r = p.add_run()
                r.text = seg
                f = r.font
                f.name = opts.get("font", FONT)
                f.size = Pt(opts.get("size", 18))
                f.bold = opts.get("bold", False)
                f.italic = opts.get("italic", False)
                f.color.rgb = opts.get("color", DARK)
            p.level = level
        return tb

    def _header(self, slide, title, accent=BLUE):
        # accent bar + title
        self._box(slide, 0, 0, W, Inches(1.05), fill=accent)
        self._text(slide, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.8),
                   [[(title, {"size": 28, "bold": True, "color": WHITE})]],
                   anchor=MSO_ANCHOR.MIDDLE)

    def _footer(self, slide, tag=""):
        self._text(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
                   [[(tag, {"size": 10, "color": GREY})]])

    # ---- slide templates -------------------------------------------------
    def _title_slide(self, title, subtitle):
        s = self._slide()
        self._box(s, 0, 0, W, H, fill=DARK)
        self._box(s, 0, Inches(2.6), W, Inches(0.08), fill=ORANGE)
        self._text(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.2),
                   [[(title, {"size": 40, "bold": True, "color": WHITE})]])
        self._text(s, Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.0),
                   [[(subtitle, {"size": 22, "color": RGBColor(0xCC, 0xCC, 0xCC)})]])
        self._text(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.6),
                   [[("Antitrust Methods: Research and Practice  ·  "
                      "Practitioner short course", {"size": 14, "color": ORANGE})]])

    def section(self, title, sub=""):
        s = self._slide()
        self._box(s, 0, 0, W, H, fill=BLUE)
        self._box(s, Inches(0.8), Inches(3.5), Inches(4), Inches(0.06), fill=WHITE)
        self._text(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.5),
                   [[(title, {"size": 34, "bold": True, "color": WHITE})]],
                   anchor=MSO_ANCHOR.BOTTOM)
        if sub:
            self._text(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.2),
                       [[(sub, {"size": 18, "color": RGBColor(0xDD, 0xEE, 0xFF)})]])
        return s

    def bullets(self, title, items, accent=BLUE, footer="", sub=""):
        """items: list of strings or (text, level) tuples."""
        s = self._slide()
        self._header(s, title, accent)
        top = Inches(1.3)
        runs = []
        if sub:
            runs.append([(sub, {"size": 16, "italic": True, "color": GREY})])
        for it in items:
            if isinstance(it, tuple):
                txt, lvl = it
            else:
                txt, lvl = it, 0
            bullet = "•  " if lvl == 0 else "–  "
            size = 20 if lvl == 0 else 17
            color = DARK if lvl == 0 else GREY
            bold = False
            runs.append([(bullet + txt, {"size": size, "color": color,
                                         "bold": bold, "level": lvl})])
        self._text(s, Inches(0.6), top, Inches(12.1), Inches(5.5), runs,
                   space_after=9, line_spacing=1.05)
        if footer:
            self._footer(s, footer)
        return s

    def two_col(self, title, left_head, left_items, right_head, right_items,
                accent=BLUE, lcolor=BLUE, rcolor=GREEN, footer=""):
        s = self._slide()
        self._header(s, title, accent)
        for x, head, items, col in [
            (Inches(0.5), left_head, left_items, lcolor),
            (Inches(6.9), right_head, right_items, rcolor)]:
            self._box(s, x, Inches(1.3), Inches(5.9), Inches(0.55), fill=col)
            self._text(s, x + Inches(0.15), Inches(1.34), Inches(5.6), Inches(0.5),
                       [[(head, {"size": 18, "bold": True, "color": WHITE})]],
                       anchor=MSO_ANCHOR.MIDDLE)
            runs = []
            for it in items:
                if isinstance(it, tuple):
                    txt, lvl = it
                else:
                    txt, lvl = it, 0
                bullet = "•  " if lvl == 0 else "–  "
                runs.append([(bullet + txt, {"size": 16 if lvl == 0 else 14,
                             "color": DARK if lvl == 0 else GREY, "level": lvl})])
            self._text(s, x + Inches(0.1), Inches(2.0), Inches(5.7), Inches(4.6),
                       runs, space_after=7, line_spacing=1.03)
        if footer:
            self._footer(s, footer)
        return s

    def table(self, title, headers, rows, accent=BLUE, col_widths=None,
              footer="", note="", font_size=13):
        s = self._slide()
        self._header(s, title, accent)
        nrows, ncols = len(rows) + 1, len(headers)
        left, top = Inches(0.5), Inches(1.35)
        width = Inches(12.3)
        height = Inches(5.0 if not note else 4.6)
        gtbl = s.shapes.add_table(nrows, ncols, left, top, width, height).table
        if col_widths:
            for i, cw in enumerate(col_widths):
                gtbl.columns[i].width = Inches(cw)
        # header
        for j, htext in enumerate(headers):
            c = gtbl.cell(0, j)
            c.fill.solid(); c.fill.fore_color.rgb = accent
            self._cell(c, htext, WHITE, True, font_size + 1)
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                c = gtbl.cell(i, j)
                c.fill.solid()
                c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
                self._cell(c, str(val), DARK, False, font_size)
        if note:
            self._text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7),
                       [[(note, {"size": 12, "italic": True, "color": GREY})]])
        if footer:
            self._footer(s, footer)
        return s

    def _cell(self, cell, text, color, bold, size):
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = FONT

    def formula(self, title, intro, formula_lines, explain, accent=ORANGE,
                footer=""):
        """A highlighted formula/definition slide."""
        s = self._slide()
        self._header(s, title, accent)
        if intro:
            self._text(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.9),
                       [[(intro, {"size": 18, "color": DARK})]],
                       line_spacing=1.05)
        # formula panel
        self._box(s, Inches(1.2), Inches(2.45), Inches(10.9), Inches(1.5),
                  fill=LIGHT, line=ORANGE)
        fruns = [[(ln, {"size": 24, "bold": True, "font": MONO, "color": BLUE})]
                 for ln in formula_lines]
        self._text(s, Inches(1.4), Inches(2.6), Inches(10.5), Inches(1.2),
                   fruns, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                   line_spacing=1.1)
        runs = []
        for it in explain:
            runs.append([("•  " + it, {"size": 17, "color": DARK})])
        self._text(s, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.4),
                   runs, space_after=8, line_spacing=1.05)
        if footer:
            self._footer(s, footer)
        return s

    def case_box(self, title, jurisdiction, holding, facts, lesson,
                 accent=None, footer=""):
        """A case-study slide with jurisdiction tag."""
        col = JURIS.get(jurisdiction, GREY)
        accent = accent or col
        s = self._slide()
        self._header(s, title, accent)
        # jurisdiction chip
        self._box(s, Inches(0.6), Inches(1.3), Inches(1.6), Inches(0.5), fill=col)
        self._text(s, Inches(0.6), Inches(1.32), Inches(1.6), Inches(0.46),
                   [[(jurisdiction, {"size": 18, "bold": True, "color": WHITE})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._text(s, Inches(2.4), Inches(1.3), Inches(10.3), Inches(0.6),
                   [[(holding, {"size": 18, "bold": True, "color": col})]],
                   anchor=MSO_ANCHOR.MIDDLE)
        # facts
        runs = [[("Facts & analysis", {"size": 15, "bold": True, "color": GREY})]]
        for f in facts:
            runs.append([("•  " + f, {"size": 16, "color": DARK})])
        self._text(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(3.4),
                   runs, space_after=7, line_spacing=1.04)
        # lesson panel
        self._box(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.1),
                  fill=LIGHT, line=accent)
        self._text(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.95),
                   [[("Why it matters:  ", {"size": 16, "bold": True, "color": accent}),
                     (lesson, {"size": 16, "color": DARK})]],
                   anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.03)
        if footer:
            self._footer(s, footer)
        return s

    def code_demo(self, title, intro, code, takeaway, accent=GREEN, footer=""):
        """A 'code as demo' slide: shows R snippet + what to point at."""
        s = self._slide()
        self._header(s, title, accent)
        if intro:
            self._text(s, Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.7),
                       [[(intro, {"size": 16, "italic": True, "color": GREY})]])
        self._box(s, Inches(0.6), Inches(1.95), Inches(12.1), Inches(3.4),
                  fill=RGBColor(0x1E, 0x1E, 0x1E))
        cruns = [[(ln, {"size": 14, "font": MONO,
                        "color": RGBColor(0xE6, 0xE6, 0xE6)})] for ln in code]
        self._text(s, Inches(0.85), Inches(2.1), Inches(11.7), Inches(3.1),
                   cruns, space_after=2, line_spacing=1.0)
        self._box(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(1.2),
                  fill=LIGHT, line=accent)
        self._text(s, Inches(0.8), Inches(5.65), Inches(11.7), Inches(1.0),
                   [[("Point at:  ", {"size": 16, "bold": True, "color": accent}),
                     (takeaway, {"size": 16, "color": DARK})]],
                   anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.03)
        if footer:
            self._footer(s, footer)
        return s

    def exercise(self, title, prompt_lines, deliverable="", accent=PURPLE,
                 footer=""):
        s = self._slide()
        self._header(s, title, accent)
        self._box(s, Inches(0.6), Inches(1.3), Inches(0.12), Inches(4.0),
                  fill=accent)
        runs = []
        for ln in prompt_lines:
            if isinstance(ln, tuple):
                txt, lvl = ln
            else:
                txt, lvl = ln, 0
            runs.append([(("" if lvl else "") + txt,
                          {"size": 19 if lvl == 0 else 16,
                           "color": DARK if lvl == 0 else GREY,
                           "bold": lvl == 0 and txt.endswith("?") is False and False,
                           "level": lvl})])
        self._text(s, Inches(1.0), Inches(1.45), Inches(11.7), Inches(4.0),
                   runs, space_after=10, line_spacing=1.1)
        if deliverable:
            self._box(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.1),
                      fill=accent)
            self._text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.95),
                       [[("Deliverable:  ", {"size": 16, "bold": True, "color": WHITE}),
                         (deliverable, {"size": 16, "color": WHITE})]],
                       anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.03)
        if footer:
            self._footer(s, footer)
        return s

    def takeaways(self, title, items, accent=DARK, footer=""):
        s = self._slide()
        self._header(s, title, accent)
        runs = []
        for i, it in enumerate(items, 1):
            runs.append([(f"{i}.  ", {"size": 20, "bold": True, "color": ORANGE}),
                         (it, {"size": 20, "color": DARK})])
        self._text(s, Inches(0.7), Inches(1.4), Inches(12.0), Inches(5.4),
                   runs, space_after=14, line_spacing=1.08)
        if footer:
            self._footer(s, footer)
        return s

    def save(self, path):
        self.prs.save(path)
        return path
